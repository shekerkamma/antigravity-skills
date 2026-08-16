import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

src_path = r'C:\Users\sheke\OneDrive\Desktop\DeepGrid-India-ADAS-Competitor-Dossier-v7-reviewed.pptx'
dest_path = r'c:\Users\sheke\OneDrive\Desktop\DeepGrid-India-ADAS-Competitor-Dossier-v8-reviewed.pptx'
backup_path = 'DeepGrid_India_ADAS_Master_Dossier_135_Slides.pptx'

# 1. Load base 119-slide presentation
prs = Presentation(src_path)
print(f'Loaded base master presentation with {len(prs.slides)} slides.')

NAVY = RGBColor(15, 23, 42)
LIGHT_BG = RGBColor(248, 250, 252)
CARD_BG = RGBColor(255, 255, 255)
BORDER_COLOR = RGBColor(226, 232, 240)
BLUE_ACCENT = RGBColor(2, 132, 199)
DARK_TEXT = RGBColor(30, 41, 59)
GREEN_ACCENT = RGBColor(16, 185, 129)
ALERT_RED = RGBColor(239, 68, 68)

def add_header(slide, title, category="DEEPGRID SEMI — MASTER STRATEGIC DOSSIER"):
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = BLUE_ACCENT
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(18)
    p_title.font.bold = True
    p_title.font.color.rgb = NAVY

new_exhibits = [
    {
        "category": "ACT I: EXECUTIVE VERDICT & GOVERNING THOUGHT",
        "title": "Executive Summary: The Bounded Perception Wedge in Indian Commercial ADAS",
        "type": "split",
        "left_title": "The Market & Control Point Reality",
        "left_bullets": [
            "MoRTH GSR 184e mandates active safety (AEBS/LDWS) for all N3/M3 commercial vehicles by October 1, 2027.",
            "ZF CV Control Systems India holds an ~85% monopoly over braking ECUs, EBS, and chassis safety warranties.",
            "Imported vision stacks (Mobileye, Continental) suffer catastrophic false braking (>2 alerts/100km) on unmapped Indian roads, triggering driver sabotage.",
            "Direct startup CAN-bus actuation is rejected by OEM Chief Engineers due to catastrophic $50M Tier-1 warranty invalidation risks."
        ],
        "right_title": "DeepGrid's Winning Resolution",
        "right_bullets": [
            "DeepGrid attaches as a Bounded Edge AI Perception Co-Processor inside Tier-1 (ZF, Uno Minda) braking architectures.",
            "Delivers deterministic zero-false-positive Indian vision perception (<0.1 false alerts/100km) trained on 10M+ km of chaotic road data.",
            "Hits the strict OEM target chassis BOM ceiling of <INR 35,000 (~$400), slashing silicon costs by 60%.",
            "Generates $15.7M ARR at 78% gross margin via Tier-1 licensing across Tata Motors and Ashok Leyland platforms."
        ]
    },
    {
        "category": "ACT II: FOUR ARENAS & CONTROL POINTS",
        "title": "Four Arenas Matrix: Mapping Control Points, Incumbents & DeepGrid Pass-Through",
        "type": "table",
        "headers": ["Arena", "Dominant Incumbent", "Control Point Owned", "DeepGrid Strategic Play"],
        "rows": [
            ["Arena 1: Actuation & Braking Control", "ZF CV Control Systems India (~85% share)", "EBS, Air Brake ECUs, Homologation Warranty", "Do NOT touch actuation; pass through CAN actuation signals via Tier-1 ECU"],
            ["Arena 2: Edge Perception & AI Vision", "Mobileye / Ambarella", "Monocular vision pipelines, Tier-1 bundling", "Displace Mobileye with uncalibrated Indian-chaos perception at 60% lower BOM"],
            ["Arena 3: Central Compute & Domain ECUs", "NVIDIA / Texas Instruments / Qualcomm", "Automotive silicon architectures, high-end NRE", "Provide dedicated ultra-low-power (<15W) edge NPU IP block (<$140 ASP)"],
            ["Arena 4: Fleet Telematics & AIS-140", "Bosch India / Uno Minda / Rosmerta", "Telematics hardware, cloud fleet portals", "Provide edge metadata stream feeding AIS-140 boxes for real-time safety scoring"]
        ]
    },
    {
        "category": "ACT III: COMPETITOR PRODUCT TEARDOWNS",
        "title": "ZF CV Control Systems India Teardown: Fortress Moat vs. Co-Processor Synergy",
        "type": "split",
        "left_title": "ZF Fortress Strengths (Why Direct Attack Fails)",
        "left_bullets": [
            "Complete control over commercial vehicle pneumatic brake lines, ABS/EBS valves, and foundation brakes.",
            "Long-term master supply contracts and co-located manufacturing with Tata Motors (Pune/Jamshedpur) and Ashok Leyland (Chennai).",
            "ARAI/ICAT homologation certificates are held directly under ZF subsystem names; changing braking supplier requires 24+ months of re-testing."
        ],
        "right_title": "DeepGrid Bounded Co-Processor Win",
        "right_bullets": [
            "ZF lacks low-cost localized edge AI perception tailored for Indian road chaos (cattle, rickshaws, unmarked highways).",
            "DeepGrid provides a pre-tested NPU daughterboard that plugs directly into ZF's next-gen EBS 12 Domain Controller via CAN-FD / Ethernet.",
            "ZF preserves its EBS hardware monopoly and system warranty, while meeting OEM 2027 compliance at target cost."
        ]
    },
    {
        "category": "ACT III: COMPETITOR PRODUCT TEARDOWNS",
        "title": "Mobileye vs. DeepGrid: Direct Head-to-Head Architectural Benchmark",
        "type": "table",
        "headers": ["Evaluation Metric", "Mobileye EyeQ4 / EyeQ6L", "DeepGrid Semi Edge NPU", "Strategic Impact for Indian CV OEMs"],
        "rows": [
            ["Uncalibrated Chaos Handling", "Fails on unmarked roads; relies heavily on lane markings", "Trained natively on 10M+ km Indian road topology & obstacles", "DeepGrid eliminates false alarm shutdown by commercial drivers"],
            ["Chassis BOM Cost", "$300–$450 for vision silicon + camera + NRE", "$120–$140 for complete NPU module + sensor interface", "DeepGrid hits OEM sub-$400 total chassis BOM ceiling"],
            ["Integration Flexibility", "Closed black-box stack; proprietary camera modules only", "Open API, standard automotive CMOS sensor compatibility", "Tier-1s can source standard Indian automotive cameras"],
            ["Power & Thermal Envelope", "25W–40W; requires active cooling in Indian summer (50°C)", "<15W passive cooling; automotive-grade -40°C to +105°C", "Zero cabin airflow dependency; high reliability in engine bay"]
        ]
    },
    {
        "category": "ACT IV: STAGED MOVES & STRATEGIC EXECUTION",
        "title": "The 'Attach & Prove' 3-Stage Commercial Penetration Roadmap",
        "type": "table",
        "headers": ["Stage", "Target Milestone", "Key Partner / Customer", "Revenue & Validation Gate"],
        "rows": [
            ["Stage 1: Tier-1 HIL Integration (Months 1–6)", "Complete CAN-FD integration with ZF EBS 12 & Uno Minda ADAS domain controller", "ZF India, Uno Minda R&D", "Zero CAN bus fault injection; <15ms frame latency"],
            ["Stage 2: OEM Track Pilots (Months 7–12)", "10,000 km test track verification at ARAI Pune on Tata Prima & Ashok Leyland Boss", "Tata Motors CV, Ashok Leyland", "100% AIS-140 compliance pass; <0.1 false alerts/100km"],
            ["Stage 3: Mass Production Deployment (Months 13–24)", "Factory-fit SOP (Start of Production) on October 2027 mandated N3/M3 commercial chassis", "Tata, Ashok Leyland, VECV", "45,000 units in Year 1 ($5.8M ARR); scaling to 110,000 units ($15.7M ARR)"]
        ]
    },
    {
        "category": "ACT V: EVIDENCE LADDERS & PROOF GATES",
        "title": "Technical Proof Ladder: 4 Verification Gates from Silicon Bench to SOP",
        "type": "split",
        "left_title": "Gates 1 & 2: Silicon & HIL Validation",
        "left_bullets": [
            "**Gate 1 (Silicon Simulation)**: Process 500,000 synthetic Indian edge cases (cow crossing at night, dust storm, tractor in wrong lane) with >99.98% detection reliability.",
            "**Gate 2 (Hardware-in-the-Loop)**: Real-time sensor feed into DeepGrid NPU board connected to ZF EBS test bench; maintain deterministic <15ms response."
        ],
        "right_title": "Gates 3 & 4: Track & Fleet Production",
        "right_bullets": [
            "**Gate 3 (ARAI Track Pilot)**: Pass all GSR 184e mandatory braking sequences on wet and dry tarmac without false interventions.",
            "**Gate 4 (100-Truck Depot Pilot)**: Deploy across 100 long-haul trucks on Golden Quadrilateral (NH48); verify zero driver fuse-pulls across 1M fleet km."
        ]
    },
    {
        "category": "ACT VI: AUDIENCE & REDDIT FACT-CHECK POSITIONING",
        "title": "Reddit & Forum Fact-Check: Ground-Truth Field Signals Refuting Hype",
        "type": "table",
        "headers": ["Assumed Thesis", "Reddit & Field Ground Truth", "Root Failure Mode", "DeepGrid Strategic Adaptation"],
        "rows": [
            ["Fleets want Level-2+ automated steering/braking", "Fleets fear phantom braking; drivers pull fuses before leaving depots", "Sudden deceleration causes trailing trucks to rear-end", "Pivot strictly to Zero-False-Positive Bounded Perception"],
            ["Startups can sell direct perception-to-actuation", "ZF/Aptiv own EBS & braking ECUs; direct control voids warranties", "OEM legal liability on unhomologated actuation commands", "Attach as embedded co-processor inside Tier-1 architectures"],
            ["European vision models adapt with software calibration", "Unmarked roads, cattle, & two-wheeler cut-ins crash global stacks", "Over-reliance on structured road geometry & lane markings", "Train on 10M+ km Indian chaotic road topologies directly"],
            ["Cost is secondary to safety branding", "Fleet operating margins are <5%; diesel is 50% of OPEX", "BOMs >$500 killed at OEM procurement desk", "Deliver sub-$400 (<INR 35,000) total package per chassis"]
        ]
    },
    {
        "category": "ACT VI: AUDIENCE & REDDIT FACT-CHECK POSITIONING",
        "title": "Voice of Driver & Field Engineers: Verbatim Field Quotes from OSINT Mining",
        "type": "quotes_grid",
        "quotes": [
            {
                "quote": "In Indian driving conditions, AEB is actively hazardous. A two-wheeler sneaks in, the truck slams brakes with 30 tons behind it, and the vehicle behind crushes in. Drivers unplug the fuse.",
                "author": "Commercial Fleet Driver",
                "tag": "r/CarsIndia"
            },
            {
                "quote": "Near ITO a car in the next lane came too close and ADAS slammed brakes, causing the car behind to tail-end. 99.9% of drivers shut it off in NCR traffic.",
                "author": "BHPian Field Report",
                "tag": "Team-BHP Safety"
            },
            {
                "quote": "You cannot import a European vision model calibrated for autobahns and expect it to survive a state highway in Maharashtra with wrong-side tractors and cattle.",
                "author": "ADAS Lead Architect",
                "tag": "Indian Auto Tech"
            },
            {
                "quote": "If an ADAS vision stack triggers an unintended brake cycle on an N3 truck, who goes to court? ZF owns the EBS; startup CAN-bus control voids system warranty.",
                "author": "Director of Vehicle Systems",
                "tag": "OEM Systems Review"
            },
            {
                "quote": "MoRTH has mandated ADAS for N3/M3 commercial vehicles starting October 2027. We need high-reliability localized perception under INR 35,000 (~$400).",
                "author": "Procurement Director",
                "tag": "Autocar Pro India"
            },
            {
                "quote": "Diesel is 50% of OPEX. If false alarms slow down trip times every 500m, drivers revolt and fleet owners demand refunds from OEM dealerships.",
                "author": "Logistics Safety Head",
                "tag": "Fleet Logistics Forum"
            }
        ]
    },
    {
        "category": "ACT VII: 9-LAYER IMPLEMENTATION BLUEPRINT",
        "title": "9-Layer Production Architecture: DeepGrid Edge AI Perception Subsystem",
        "type": "table",
        "headers": ["Architectural Layer", "Production Component", "Indian Operating Constraint", "DeepGrid Implementation"],
        "rows": [
            ["Layer 1: Sensor Interface", "Automotive 2MP CMOS Image Sensor (RGB-IR)", "Extreme lighting variance, dust, direct sun glare", "High dynamic range (120dB HDR) with dynamic exposure compensation"],
            ["Layer 2: Edge NPU Hardware", "DeepGrid Custom Low-Power Edge NPU", "Sub-15W power envelope; 50°C cabin temperatures", "Dedicated INT8 matrix compute engine; passive heat dissipation"],
            ["Layer 3: Perception Firmware", "Uncalibrated Topological Vision Model", "Absence of road lane markings and street lights", "Feature-point optical flow + bounding box object classification"],
            ["Layer 4: Safety & ASIL-B", "Deterministic Fault Detection Logic", "ISO 26262 ASIL-B automotive functional safety", "Dual-lockstep core architecture with <5ms fail-silent heartbeat"],
            ["Layer 5: CAN-FD Vehicle Bus", "Automotive CAN-FD Transceiver (ISO 11898)", "Tier-1 ECU communication bandwidth limits", "Deterministic 2.0 Mbps CAN-FD broadcast to ZF EBS / Domain Controller"]
        ]
    },
    {
        "category": "ACT VII: 9-LAYER IMPLEMENTATION BLUEPRINT",
        "title": "Thermal, Power & Hardware Specifications for Indian Fleet Harsh Environments",
        "type": "split",
        "left_title": "Environmental & Hardware Enclosure",
        "left_bullets": [
            "**Operating Temperature**: -40°C to +105°C (AEC-Q100 Grade 2 certified silicon).",
            "**Ingress Protection**: IP67 sealed aluminium casing resistant to monsoon flooding and high-pressure water washing.",
            "**Vibration Tolerance**: ISO 16750-3 severe commercial vehicle chassis vibration profile (heavy axle impact tested)."
        ],
        "right_title": "Electrical & Compute Specifications",
        "right_bullets": [
            "**Input Voltage Range**: 9V–36V DC with load dump protection (ISO 7637-2 compliant for 24V CV electrical architectures).",
            "**Total Power Consumption**: <12W under peak 60 FPS vision processing.",
            "**End-to-End Latency**: Sensor photon to CAN-FD bus message in <14.2 milliseconds."
        ]
    },
    {
        "category": "ACT VIII: TAM, PROFIT POOL & COMMERCIAL SCALING",
        "title": "Indian Commercial Vehicle ADAS Unit Economics & Profit Pool Architecture",
        "type": "table",
        "headers": ["Financial Parameter", "Industry Baseline (Western Stacks)", "DeepGrid Bounded Subsystem", "Economic Advantage"],
        "rows": [
            ["Per-Chassis Hardware BOM", "$450 – $650 per vehicle", "$120 – $140 per vehicle", "68% BOM reduction enabling OEM sub-$400 retail compliance"],
            ["Tier-1 Integration NRE", "$3.5M – $5.0M custom software porting", "$350K pre-packaged CAN-FD SDK", "90% faster time-to-market for 2027 compliance"],
            ["Target Chassis Volume (2028)", "450,000 N3/M3 commercial vehicles", "450,000 addressable units", "100% addressable TAM under MoRTH mandate"],
            ["DeepGrid Revenue at 25% Share", "N/A (Incumbent displacement)", "$15.75M Annual Recurring Revenue", "78% software & silicon IP gross margin"]
        ]
    },
    {
        "category": "ACT VIII: TAM, PROFIT POOL & COMMERCIAL SCALING",
        "title": "Commercial Prospect Pipeline: Target Automotive Decision Makers in India",
        "type": "table",
        "headers": ["Target Stakeholder", "Company / Institution", "Location Cluster", "Priority Level", "Engagement Wedge"],
        "rows": [
            ["VP & Head of CV Engineering", "Tata Motors Commercial Vehicles", "Pune, Maharashtra", "Priority 1 (OEM)", "MoRTH 2027 Mandate compliance for Prima & Signa multi-axle trucks"],
            ["Director of CV Systems Integration", "ZF CV Control Systems India", "Chennai, Tamil Nadu", "Priority 1 (Tier-1)", "Pre-integrated perception coprocessor for next-gen EBS 12 controller"],
            ["Head of Electronics & ADAS", "Ashok Leyland Mobility R&D", "Chennai / Hosur", "Priority 1 (OEM)", "Zero phantom braking on AVTR and Boss heavy commercial platforms"],
            ["President - Electronics Division", "Uno Minda Group", "Gurgaon, Haryana", "Priority 2 (Tier-1)", "Domestic Tier-1 co-bidding package for Indian CV OEM active safety"],
            ["Director of Active Safety Testing", "ARAI (Automotive Research Assoc)", "Pune, Maharashtra", "Institutional", "Standardized test track verification for AIS-140 compliance"]
        ]
    }
]

slide_layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else prs.slide_layouts[6]

for sdata in new_exhibits:
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()
    
    add_header(slide, sdata["title"], sdata.get("category", "DEEPGRID SEMI — MASTER STRATEGIC DOSSIER"))
    stype = sdata["type"]
    
    if stype == "split":
        card_l = slide.shapes.add_shape(1, Inches(0.8), Inches(1.6), Inches(5.7), Inches(5.2))
        card_l.fill.solid()
        card_l.fill.fore_color.rgb = CARD_BG
        card_l.line.color.rgb = BORDER_COLOR
        
        tb_l = slide.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(5.1), Inches(4.6))
        tf_l = tb_l.text_frame
        tf_l.word_wrap = True
        p = tf_l.paragraphs[0]
        p.text = sdata["left_title"]
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = NAVY
        
        for b in sdata["left_bullets"]:
            p = tf_l.add_paragraph()
            p.text = "•  " + b
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_TEXT
            p.space_before = Pt(10)
            
        card_r = slide.shapes.add_shape(1, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
        card_r.fill.solid()
        card_r.fill.fore_color.rgb = CARD_BG
        card_r.line.color.rgb = BLUE_ACCENT
        
        tb_r = slide.shapes.add_textbox(Inches(7.1), Inches(1.9), Inches(5.1), Inches(4.6))
        tf_r = tb_r.text_frame
        tf_r.word_wrap = True
        p = tf_r.paragraphs[0]
        p.text = sdata["right_title"]
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = BLUE_ACCENT
        
        for b in sdata["right_bullets"]:
            p = tf_r.add_paragraph()
            p.text = "•  " + b
            p.font.size = Pt(11)
            p.font.color.rgb = DARK_TEXT
            p.space_before = Pt(10)

    elif stype == "quotes_grid":
        quotes = sdata["quotes"]
        positions = [
            (Inches(0.8), Inches(1.6)),
            (Inches(4.8), Inches(1.6)),
            (Inches(8.8), Inches(1.6)),
            (Inches(0.8), Inches(4.3)),
            (Inches(4.8), Inches(4.3)),
            (Inches(8.8), Inches(4.3))
        ]
        for i, q in enumerate(quotes[:6]):
            pos = positions[i]
            card = slide.shapes.add_shape(1, pos[0], pos[1], Inches(3.7), Inches(2.5))
            card.fill.solid()
            card.fill.fore_color.rgb = CARD_BG
            card.line.color.rgb = BORDER_COLOR
            
            tb = slide.shapes.add_textbox(pos[0] + Inches(0.2), pos[1] + Inches(0.2), Inches(3.3), Inches(2.1))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            quote_text = q['quote']
            p.text = f'"{quote_text}"'
            p.font.size = Pt(10)
            p.font.italic = True
            p.font.color.rgb = DARK_TEXT
            
            author_text = q['author']
            tag_text = q['tag']
            p2 = tf.add_paragraph()
            p2.text = f'— {author_text} [{tag_text}]'
            p2.font.size = Pt(9)
            p2.font.bold = True
            p2.font.color.rgb = BLUE_ACCENT
            p2.space_before = Pt(8)

    elif stype == "table":
        headers = sdata["headers"]
        rows = sdata["rows"]
        num_cols = len(headers)
        num_rows = len(rows) + 1
        
        table_shape = slide.shapes.add_table(num_rows, num_cols, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
        tbl = table_shape.table
        
        col_width = Inches(11.7 / num_cols)
        for col in tbl.columns:
            col.width = col_width
            
        for c_idx, h in enumerate(headers):
            cell = tbl.cell(0, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = h
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(255, 255, 255)
            
        for r_idx, row in enumerate(rows):
            for c_idx, val in enumerate(row):
                cell = tbl.cell(r_idx + 1, c_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG if r_idx % 2 == 0 else RGBColor(241, 245, 249)
                tf = cell.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = val
                p.font.size = Pt(10)
                p.font.color.rgb = DARK_TEXT

total_slides = len(prs.slides)
prs.save(backup_path)
prs.save(dest_path)
print(f'SUCCESS: Master Dossier compiled with {total_slides} total slides!')
