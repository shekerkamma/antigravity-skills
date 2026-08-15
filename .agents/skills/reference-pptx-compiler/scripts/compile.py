import sys
import json
import re
from pptx import Presentation
from pptx.enum.text import PP_ALIGN

def populate_table_shape(table_shape, headers, rows):
    table = table_shape.table
    
    # Check dimensions
    req_rows = len(rows) + (1 if headers else 0)
    req_cols = len(headers) if headers else (len(rows[0]) if rows else 1)
    
    # We write headers in row 0
    current_row = 0
    if headers:
        for col_idx, h_text in enumerate(headers):
            if col_idx < len(table.columns):
                cell = table.cell(0, col_idx)
                cell.text = str(h_text)
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.LEFT
                    p.font.bold = True
        current_row += 1
        
    # Write data rows
    for r_idx, row in enumerate(rows):
        target_r = current_row + r_idx
        if target_r < len(table.rows):
            for c_idx, val in enumerate(row):
                if c_idx < len(table.columns):
                    cell = table.cell(target_r, c_idx)
                    cell.text = str(val)
                    for p in cell.text_frame.paragraphs:
                        p.alignment = PP_ALIGN.LEFT

def compile_custom_deck(findings_path, template_path, layout_map_path, output_path):
    # Load dynamic content payload
    with open(findings_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
        
    # Load layout mappings
    with open(layout_map_path, 'r', encoding='utf-8') as f:
        layout_map = json.load(f)
        
    prs = Presentation(template_path)
    
    company = data.get('company', '').strip()
    company_upper = company.upper()
    
    slides_to_keep = set()
    
    # 1. Map and Populate Title Slide (index 0 / domain TITLE)
    title_slide_idx = 0
    for slide_info in layout_map.get('slides', []):
        if slide_info.get('domain') == 'TITLE':
            title_slide_idx = slide_info['slide_index']
            break
            
    slides_to_keep.add(title_slide_idx)
    title_slide = prs.slides[title_slide_idx]
    
    title_map = next((s for s in layout_map['slides'] if s['slide_index'] == title_slide_idx), None)
    if title_map:
        t_mappings = title_map.get('mappings', {})
        h_idx = t_mappings.get('headline')
        r_idx = t_mappings.get('executive_read')
        
        if h_idx is not None and h_idx < len(title_slide.shapes):
            title_slide.shapes[h_idx].text = data.get('headline', '')
        if r_idx is not None and r_idx < len(title_slide.shapes):
            title_slide.shapes[r_idx].text = f"Template: superdesign-theme  ·  {company}  ·  June 2026"

    # 2. Map and Populate Dynamic Slide Payloads
    findings_slides = data.get('slides', [])
    
    for s_payload in findings_slides:
        s_domain = s_payload.get('domain', '').upper().strip()
        
        # Locate corresponding template slide index
        target_slide_idx = -1
        for slide_info in layout_map.get('slides', []):
            if slide_info.get('domain') == s_domain:
                target_slide_idx = slide_info['slide_index']
                break
                
        if target_slide_idx == -1:
            print(f"Warning: Layout for slide domain '{s_domain}' not found in template. Skipping.")
            continue
            
        slides_to_keep.add(target_slide_idx)
        target_slide = prs.slides[target_slide_idx]
        
        # Get target mappings
        target_map = next((s for s in layout_map['slides'] if s['slide_index'] == target_slide_idx), None)
        if not target_map:
            continue
            
        mappings = target_map.get('mappings', {})
        
        # Inject Headline
        h_idx = mappings.get('headline')
        if h_idx is not None and h_idx < len(target_slide.shapes) and 'headline' in s_payload:
            target_slide.shapes[h_idx].text = s_payload['headline']
            
        # Inject Executive Read
        r_idx = mappings.get('executive_read')
        if r_idx is not None and r_idx < len(target_slide.shapes) and 'executive_read' in s_payload:
            target_slide.shapes[r_idx].text = s_payload['executive_read']
            
        # Inject Cards
        card_configs = mappings.get('cards', [])
        payload_bullets = s_payload.get('bullets', [])
        
        for k, c_config in enumerate(card_configs):
            if k < len(payload_bullets):
                bullet_item = payload_bullets[k]
                title_idx = c_config.get('title_idx')
                body_idx = c_config.get('body_idx')
                prefix = c_config.get('prefix', "")
                
                # Check if item is structured as "Title: Body" or simple text
                if isinstance(bullet_item, str) and ":" in bullet_item:
                    parts = bullet_item.split(":", 1)
                    card_title = parts[0].strip().upper()
                    card_body = "• " + parts[1].strip()
                else:
                    card_title = f"POINT {k+1}"
                    card_body = "• " + str(bullet_item).strip()
                    
                if title_idx is not None and title_idx < len(target_slide.shapes):
                    target_slide.shapes[title_idx].text = f"{prefix}{card_title}"
                if body_idx is not None and body_idx < len(target_slide.shapes):
                    target_slide.shapes[body_idx].text = card_body
            else:
                # Clear unused card
                title_idx = c_config.get('title_idx')
                body_idx = c_config.get('body_idx')
                if title_idx is not None and title_idx < len(target_slide.shapes):
                    target_slide.shapes[title_idx].text = ""
                if body_idx is not None and body_idx < len(target_slide.shapes):
                    target_slide.shapes[body_idx].text = ""
                    
        # Inject Tables
        t_config = mappings.get('table')
        if 'table_data' in s_payload:
            t_payload = s_payload['table_data']
            headers = t_payload.get('headers', [])
            rows = t_payload.get('rows', [])
            
            # Find the actual PPTX Table shape on the slide
            table_shape = None
            for shape in target_slide.shapes:
                if shape.has_table:
                    table_shape = shape
                    break
                    
            if table_shape:
                populate_table_shape(table_shape, headers, rows)
            elif t_config:
                # Fallback to individual shapes if configured in layout map
                start_shape = t_config['start_shape_idx']
                stride = t_config.get('stride', 5)
                fields = t_config.get('fields', [])
                
                current_idx = start_shape
                for row_idx, row in enumerate(rows):
                    for f_offset, field in enumerate(fields):
                        cell_idx = current_idx + f_offset
                        if cell_idx < len(target_slide.shapes) and f_offset < len(row):
                            target_slide.shapes[cell_idx].text = str(row[f_offset])
                    current_idx += stride

    # 3. Clean up presentation: keep only the populated slides
    slides_to_delete = [i for i in range(len(prs.slides)) if i not in slides_to_keep]
    for i in sorted(slides_to_delete, reverse=True):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
        
    total_slides = len(prs.slides)
    
    # 4. Global Text Replacement and Footer Re-calculation
    # Re-calculate index offsets
    kept_indices_mapping = {old_idx: new_idx for new_idx, old_idx in enumerate(sorted(list(slides_to_keep)))}
    
    for s_idx, slide in enumerate(prs.slides):
        current_slide_num = s_idx + 1
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                orig_text = shape.text
                new_text = orig_text
                
                # Replace Template Brand Name
                new_text = re.sub(r'\bSUPERDESIGN\b', company_upper, new_text, flags=re.IGNORECASE)
                new_text = re.sub(r'\bACCENTURE\b', company_upper, new_text, flags=re.IGNORECASE)
                
                # Update footer builder tag
                new_text = new_text.replace("Generated via Superdesign PPTX Builder", f"Generated via Superdesign PPTX Builder  ·  {company}")
                
                # Re-calculate slide count footers (e.g. "2 / 5" -> "2 / 4")
                new_text = re.sub(r'\b\d+\s*/\s*\d+\b', f"{current_slide_num} / {total_slides}", new_text)
                
                # Apply changes if modified
                if new_text != orig_text:
                    shape.text = new_text
                    
    prs.save(output_path)
    print(f"Successfully compiled {output_path} with {len(prs.slides)} slides.")

if __name__ == "__main__":
    if len(sys.argv) < 5:
        print("Usage: python compile.py <findings_json> <template_pptx> <layout_map_json> <output_pptx>")
        sys.exit(1)
    compile_custom_deck(sys.argv[1], sys.argv[2], sys.argv[3], sys.argv[4])
