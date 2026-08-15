import sys
import json
from pptx import Presentation

def analyze_and_map_pptx(pptx_path, output_json_path):
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Error opening presentation: {e}")
        sys.exit(1)
        
    layout_map = {
        "slides": []
    }
    
    for i, slide in enumerate(prs.slides):
        slide_info = {
            "slide_index": i,
            "slide_title": f"Slide {i+1}",
            "domain": None,
            "shapes": [],
            "mappings": {
                "headline": None,
                "executive_read": None,
                "key_findings": [],
                "table": None
            }
        }
        
        # We search shapes
        findings_shapes = []
        table_shapes = []
        
        for j, shape in enumerate(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip():
                txt = shape.text.strip()
                shape_detail = {
                    "idx": j,
                    "text": txt[:100],
                    "shape_type": str(shape.shape_type)
                }
                slide_info["shapes"].append(shape_detail)
                
                # Check for framework domain name (e.g. contains "·" or "DOMAIN")
                if "·" in txt or "DOMAIN:" in txt.upper():
                    slide_info["slide_title"] = txt
                    # Extract domain suffix
                    parts = txt.split('·')
                    slide_info["domain"] = parts[-1].strip().upper()
                
                # Try auto-mapping common placeholders
                txt_upper = txt.upper()
                if "HEADLINE" in txt_upper or "ASSERTION-STYLE CLAIM" in txt_upper:
                    slide_info["mappings"]["headline"] = j
                elif "EXECUTIVE READ" in txt_upper or "EXECUTIVE SUMMARY" in txt_upper or "2-3 SENTENCE" in txt_upper:
                    slide_info["mappings"]["executive_read"] = j
                elif "FINDING" in txt_upper or "BULLET" in txt_upper:
                    findings_shapes.append(j)
                elif "ACCENTURE OUTPUT" in txt_upper or "TABLE_ROW" in txt_upper or "DIMENSION" in txt_upper:
                    table_shapes.append(j)
                    
        # Sort and save findings shapes
        if findings_shapes:
            slide_info["mappings"]["key_findings"] = findings_shapes
            
        # Try to map table patterns if shape sequences are detected
        if table_shapes:
            # Table starts roughly around the first detected row-indicator shape
            start_shape = min(table_shapes)
            # Find the total length of cells/shapes at the end of the slide
            # e.g., if there are 4 columns (area, evidence, interpretation, confidence) the stride is 4 or 5
            # Let's inspect shapes after start_shape to see if they look like a grid
            slide_info["mappings"]["table"] = {
                "start_shape_idx": start_shape,
                "stride": 5, # Default stride of 5 for area, evidence, interpretation, confidence, spacer
                "fields": ["area", "evidence", "interpretation", "confidence"]
            }
            
        # Apply title slide heuristics if it's the first slide and no domain found
        if i == 0 and not slide_info["domain"]:
            slide_info["domain"] = "TITLE"
            # Auto-map slide 1 shapes
            title_shapes = [s["idx"] for s in slide_info["shapes"] if len(s["text"]) > 5]
            if len(title_shapes) >= 3:
                slide_info["mappings"]["headline"] = title_shapes[1]  # Typically company title
                slide_info["mappings"]["executive_read"] = title_shapes[2]  # Subtitle
        
        layout_map["slides"].append(slide_info)
        
    with open(output_json_path, 'w', encoding='utf-8') as f:
        json.dump(layout_map, f, indent=2)
        
    print(f"Successfully mapped {len(prs.slides)} slides and saved to {output_json_path}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python map.py <reference_pptx> <output_json_path>")
        sys.exit(1)
    analyze_and_map_pptx(sys.argv[1], sys.argv[2])
