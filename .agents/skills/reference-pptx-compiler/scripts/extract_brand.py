import sys
import json
import collections
from pptx import Presentation
from pptx.dml.color import RGBColor

def rgb_to_hex(rgb):
    return "#{:02x}{:02x}{:02x}".format(rgb[0], rgb[1], rgb[2])

def extract_brand_tokens(template_path, output_json_path):
    try:
        prs = Presentation(template_path)
    except Exception as e:
        print(f"Error opening template: {e}")
        sys.exit(1)
        
    bg_colors = []
    shape_fills = []
    text_colors = []
    fonts = []
    
    for slide in prs.slides:
        # 1. Background color
        try:
            if slide.background.fill.type == 1:  # Solid
                bg_colors.append(slide.background.fill.fore_color.rgb)
        except Exception:
            pass
            
        for shape in slide.shapes:
            # 2. Shape fill colors
            try:
                # If shape covers the whole slide, treat it as a background color
                is_full_slide = (shape.left <= 0 and shape.top <= 0 and 
                                 shape.width >= prs.slide_width and 
                                 shape.height >= prs.slide_height)
                
                if shape.fill.type == 1:
                    rgb = shape.fill.fore_color.rgb
                    if is_full_slide:
                        bg_colors.append(rgb)
                    else:
                        shape_fills.append(rgb)
            except Exception:
                pass
                
            # 3. Text colors & fonts
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            fonts.append(run.font.name)
                        if run.font.color and run.font.color.type == 1:
                            text_colors.append(run.font.color.rgb)
                            
    # Count frequencies
    bg_counter = collections.Counter(bg_colors)
    fill_counter = collections.Counter(shape_fills)
    text_counter = collections.Counter(text_colors)
    font_counter = collections.Counter(fonts)
    
    # 1. Background Color (default to dark navy if not found)
    bg_color = bg_counter.most_common(1)[0][0] if bg_colors else RGBColor(0x0A, 0x16, 0x28)
    
    # 2. Main Font
    main_font = font_counter.most_common(1)[0][0] if fonts else "Calibri"
    
    # 3. Accent Colors
    # Filter out colors close to background, pure black, or pure white
    filtered_fills = []
    for rgb in shape_fills:
        # Simple distance check from black/white
        is_near_white = (rgb[0] > 240 and rgb[1] > 240 and rgb[2] > 240)
        is_near_black = (rgb[0] < 25 and rgb[1] < 25 and rgb[2] < 25)
        is_near_bg = (abs(rgb[0] - bg_color[0]) < 20 and 
                      abs(rgb[1] - bg_color[1]) < 20 and 
                      abs(rgb[2] - bg_color[2]) < 20)
                      
        if not (is_near_white or is_near_black or is_near_bg):
            filtered_fills.append(rgb)
            
    fill_counter_filtered = collections.Counter(filtered_fills)
    common_accents = fill_counter_filtered.most_common(3)
    
    accent_1 = common_accents[0][0] if len(common_accents) > 0 else RGBColor(0x00, 0x9B, 0x82)
    accent_2 = common_accents[1][0] if len(common_accents) > 1 else RGBColor(0x00, 0xC9, 0xA7)
    accent_3 = common_accents[2][0] if len(common_accents) > 2 else RGBColor(0xFF, 0xB8, 0x00)
    
    # Export tokens as HEX strings
    tokens = {
        "NAVY": rgb_to_hex(bg_color),
        "ACCENT": rgb_to_hex(accent_1),
        "TEAL": rgb_to_hex(accent_2),
        "GOLD": rgb_to_hex(accent_3),
        "FONT": main_font,
        "FONT_H": main_font
    }
    
    with open(output_json_path, 'w', encoding='utf-8') as f:
        json.dump(tokens, f, indent=2)
        
    print(f"Successfully extracted style tokens from {template_path}:")
    print(json.dumps(tokens, indent=2))

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python extract_brand.py <template_pptx> <output_json_path>")
        sys.exit(1)
    extract_brand_tokens(sys.argv[1], sys.argv[2])
